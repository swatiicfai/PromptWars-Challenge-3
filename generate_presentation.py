import os
import asyncio
import edge_tts
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Data for slides
slides_data = [
    {
        "title": "Terralife 🌱",
        "content": "Track, Understand & Reduce Your Carbon Footprint\nPromptWars Challenge 3 Submission",
        "voice_script": "Hello and welcome. Today, I am thrilled to present Terralife, a premium Carbon Footprint Awareness Platform designed to help individuals understand, track, and ultimately reduce their environmental impact."
    },
    {
        "title": "The Challenge",
        "content": "- High global emissions\n- Lack of individual awareness\n- Complex carbon calculations",
        "voice_script": "We face a significant global challenge. While climate change is a massive issue, many individuals lack the tools to understand their own daily impact. Calculating personal emissions is often complex and overwhelming for the average person."
    },
    {
        "title": "The Terralife Solution",
        "content": "- Multi-step interactive calculator\n- Real-time footprint estimates\n- Actionable, personalized insights",
        "voice_script": "Enter Terralife. Our solution breaks down the complexity of carbon tracking into a beautiful, gamified experience. With an interactive, multi-step calculator, users get real-time estimates of their baseline emissions across transport, energy, diet, and lifestyle."
    },
    {
        "title": "Gamified Sustainability",
        "content": "- Daily Green Habits Tracker\n- Earn XP and Level up\n- Unlockable Eco-Badges",
        "voice_script": "But awareness is only the first step. Terralife drives action through gamification. Users can log daily green habits to earn Experience Points, track their live CO2 reductions, and compete in eco-challenges to unlock profile badges like 'Solar Sovereign' and 'Zero Waste Warrior'."
    },
    {
        "title": "Carbon Swap Simulator",
        "content": "- Visualize lifestyle changes\n- Interactive sliders for commuting, diet, & energy",
        "voice_script": "To further educate users, we built the Carbon Swap Simulator. This interactive tool allows individuals to play with sliders, like swapping red meat for plant-based meals, or driving for cycling, to instantly visualize the dramatic annual savings these simple swaps can achieve."
    },
    {
        "title": "Tech Stack & Architecture",
        "content": "- Frontend: Vanilla JS, HTML5, CSS3\n- Deployment: Docker & Google Cloud Run\n- Secure & Client-side heavy",
        "voice_script": "Under the hood, Terralife is built for maximum efficiency and security. By utilizing modern Vanilla JavaScript and CSS, the application is lightning fast. It runs securely on the client side, and is containerized and deployed using Google Cloud Run for seamless scalability."
    },
    {
        "title": "Thank You! 🌍",
        "content": "Try the live demo today.\nLet's build a greener future, together.",
        "voice_script": "Thank you for your time. Terralife proves that with the right design and technology, we can empower everyone to take meaningful climate action. We invite you to check out the live demo and join us in building a greener future."
    }
]

def create_pptx(output_dir):
    prs = Presentation()
    for i, slide_data in enumerate(slides_data):
        # Use a layout with Title and Body
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = slide_data['title']
        content.text = slide_data['content']
        
    prs.save(os.path.join(output_dir, 'Terralife_Presentation.pptx'))

async def create_audio(output_dir):
    # Using a high-quality female AI voice from Edge TTS
    voice = "en-US-AriaNeural"
    for i, slide_data in enumerate(slides_data):
        print(f"Generating audio for slide {i+1}...")
        communicate = edge_tts.Communicate(slide_data['voice_script'], voice)
        await communicate.save(os.path.join(output_dir, f"slide_{i+1}_voice.mp3"))

async def main():
    output_dir = "Presentation_Output"
    os.makedirs(output_dir, exist_ok=True)
    
    print("Creating PPTX...")
    create_pptx(output_dir)
    
    print("Creating Audio files...")
    await create_audio(output_dir)
    print("Done! All files are in the 'Presentation_Output' folder.")

if __name__ == "__main__":
    asyncio.run(main())
